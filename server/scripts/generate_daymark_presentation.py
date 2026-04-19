from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("/Users/chenyidian/Documents/anote/Daymark_Anote_Style_Presentation.pptx")
TEMPLATE_PATH = Path("/Users/chenyidian/Documents/anote/Anote Slide Template.pptx")

BG = RGBColor(14, 22, 35)
BG_ALT = RGBColor(18, 31, 47)
WHITE = RGBColor(247, 250, 252)
MUTED = RGBColor(171, 184, 204)
CYAN = RGBColor(63, 201, 255)
BLUE = RGBColor(84, 153, 255)
TEAL = RGBColor(109, 229, 208)
LIME = RGBColor(229, 255, 86)
SOFT_PANEL = RGBColor(32, 45, 64)
ORANGE = RGBColor(255, 155, 93)


def add_box(slide, x, y, w, h, fill, radius=0.18, line=None, alpha=0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = alpha
    shape.line.color.rgb = line or fill
    if line is None:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size=24, color=WHITE, bold=False, font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, items, size=22, color=WHITE, bullet_color=LIME, gap=0.54):
    for idx, item in enumerate(items):
        top = y + Inches(gap * idx)
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, top + Inches(0.12), Inches(0.1), Inches(0.1)).fill.solid()
        dot = slide.shapes[-1]
        dot.fill.fore_color.rgb = bullet_color
        dot.line.fill.background()
        add_textbox(slide, x + Inches(0.2), top, w - Inches(0.2), Inches(0.35), item, size=size, color=color)


def add_kicker(slide, text):
    add_textbox(slide, Inches(0.7), Inches(0.32), Inches(4.2), Inches(0.3), text, size=12, color=MUTED)


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.7), Inches(0.78), Inches(11.2), Inches(1.35), title, size=28, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.72), Inches(1.7), Inches(9.8), Inches(0.65), subtitle, size=16, color=MUTED)


def add_background(slide, bright=False):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG if not bright else BG_ALT

    orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.7), Inches(-0.3), Inches(3.2), Inches(3.2))
    orb.fill.solid()
    orb.fill.fore_color.rgb = CYAN
    orb.fill.transparency = 0.87
    orb.line.fill.background()

    orb2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-1.0), Inches(5.5), Inches(4.2), Inches(4.2))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = LIME
    orb2.fill.transparency = 0.92
    orb2.line.fill.background()

    bar_colors = [CYAN, BLUE, TEAL, LIME]
    bar_x = Inches(0.72)
    for color in bar_colors:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, bar_x, Inches(6.88), Inches(2.8), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar_x += Inches(2.72)


def add_neon_strip(slide, x, y, w, h):
    colors = [CYAN, BLUE, TEAL, LIME]
    seg_w = w / len(colors)
    current_x = x
    for color in colors:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, current_x, y, seg_w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        current_x += seg_w


def add_card(slide, x, y, w, h, title, body, accent):
    card = add_box(slide, x, y, w, h, SOFT_PANEL, alpha=0.18, line=accent)
    card.line.width = Pt(1.2)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.18), w - Inches(0.4), Inches(0.32), title, size=20, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.62), w - Inches(0.42), h - Inches(0.75), body, size=14, color=MUTED)


def add_pipeline_step(slide, x, y, w, h, title, accent):
    shape = add_box(slide, x, y, w, h, accent, alpha=0.08, line=accent)
    shape.line.width = Pt(1.4)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), h - Inches(0.24), title, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_cover_graphic(slide):
    panel1 = add_box(slide, Inches(7.2), Inches(1.28), Inches(4.1), Inches(2.55), SOFT_PANEL, alpha=0.22, line=RGBColor(90, 105, 130))
    panel1.rotation = -5
    panel2 = add_box(slide, Inches(8.55), Inches(2.08), Inches(4.15), Inches(2.65), SOFT_PANEL, alpha=0.14, line=RGBColor(150, 150, 150))
    panel2.rotation = 5
    for panel in (panel1, panel2):
        x = panel.left + Inches(0.28)
        for idx in range(5):
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, panel.top + Inches(0.35 + idx * 0.32), Inches(3.1), Inches(0.08))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(80 + idx * 16, 96 + idx * 10, 126 + idx * 8)
            line.fill.transparency = 0.15
            line.line.fill.background()
        pill = add_box(slide, panel.left + Inches(0.38), panel.top + Inches(1.65), Inches(1.0), Inches(0.22), CYAN, alpha=0.0)
        pill.fill.transparency = 0.28
        pill.line.fill.background()
    orbit = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(6.0), Inches(0.7), Inches(5.9), Inches(3.6))
    orbit.line.color.rgb = CYAN
    orbit.line.width = Pt(2)
    orbit.fill.background()
    orbit.rotation = 8


def build_presentation() -> Presentation:
    prs = Presentation()
    if TEMPLATE_PATH.exists():
        template = Presentation(str(TEMPLATE_PATH))
        prs.slide_width = template.slide_width
        prs.slide_height = template.slide_height
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_kicker(slide, "ANOTE STYLE PRESENTATION / HUMAN-CENTERED AI")
    add_textbox(slide, Inches(0.8), Inches(2.05), Inches(4.8), Inches(0.6), "Daymark", size=34, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.82), Inches(2.78), Inches(4.9), Inches(0.9), "A personal operations board that helps users decide what to keep, defer, and delete.", size=18, color=WHITE)
    add_textbox(slide, Inches(0.84), Inches(3.72), Inches(5.0), Inches(0.42), "Planning + execution + review + AI conversation", size=13, color=MUTED)
    add_cover_graphic(slide)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_background(slide, bright=True)
    add_kicker(slide, "PROBLEM / WHY THIS PRODUCT EXISTS")
    add_title(slide, "Most productivity tools optimize adding more.", "Daymark is built around a different question: how can AI help people make smaller, clearer, more realistic daily decisions?")
    add_card(slide, Inches(0.75), Inches(2.3), Inches(4.0), Inches(2.25), "The user problem", "People capture tasks in one place, execute somewhere else, and reflect almost nowhere. The day becomes fragmented, overloaded, and emotionally blind.", CYAN)
    add_card(slide, Inches(4.95), Inches(2.3), Inches(3.8), Inches(2.25), "The product gap", "Most tools treat defer and delete as failure states. Daymark treats them as intelligent decisions that reduce noise and protect focus.", TEAL)
    add_box(slide, Inches(9.05), Inches(2.15), Inches(3.25), Inches(3.7), LIME, alpha=0.0)
    add_textbox(slide, Inches(9.4), Inches(2.72), Inches(2.6), Inches(0.8), "Human-centered AI", size=26, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.35), Inches(3.7), Inches(2.7), Inches(1.4), "AI should not just generate more. It should reduce friction, surface trade-offs, and stay close to real human signals.", size=15, color=BG, align=PP_ALIGN.CENTER)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_kicker(slide, "PRODUCT THESIS")
    add_title(slide, "Daymark is not a bigger to-do list.", "It is a daily operating board where planning, execution, reflection, and emotional context live in the same loop.")
    add_neon_strip(slide, Inches(0.78), Inches(2.1), Inches(5.9), Inches(0.14))
    add_card(slide, Inches(0.8), Inches(2.55), Inches(2.8), Inches(1.65), "Inbox", "Capture one-off tasks, daily habits, and weekly commitments.", CYAN)
    add_card(slide, Inches(3.85), Inches(2.55), Inches(2.8), Inches(1.65), "Today", "Generate a realistic board, then complete, defer, or delete with intent.", BLUE)
    add_card(slide, Inches(6.9), Inches(2.55), Inches(2.8), Inches(1.65), "Review", "See what actually happened across the week and month, not just what was planned.", TEAL)
    add_card(slide, Inches(9.95), Inches(2.55), Inches(2.1), Inches(1.65), "Concierge", "Ask questions, mutate tasks, and analyze behavior conversationally.", LIME)
    add_textbox(slide, Inches(0.84), Inches(4.85), Inches(11.3), Inches(1.0), "Core idea: AI is not a detached chatbot here. It is an embedded reasoning layer inside the user's real daily workflow.", size=24, color=WHITE, bold=True)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_background(slide, bright=True)
    add_kicker(slide, "END-TO-END EXPERIENCE")
    add_title(slide, "The product loop is capture -> plan -> execute -> reflect -> adapt.")
    steps = [
        ("Capture", CYAN, "Create temporary, daily, and weekly tasks."),
        ("Plan", BLUE, "Generate a realistic Today board."),
        ("Execute", TEAL, "Run focus sessions and mark outcomes."),
        ("Reflect", LIME, "Review history, patterns, and signals."),
        ("Ask AI", ORANGE, "Use natural language to query and operate the app."),
    ]
    step_x = Inches(0.8)
    for label, accent, body in steps:
        add_pipeline_step(slide, step_x, Inches(2.35), Inches(2.2), Inches(0.7), label, accent)
        add_textbox(slide, step_x, Inches(3.18), Inches(2.2), Inches(0.95), body, size=13, color=MUTED)
        step_x += Inches(2.4)
    add_textbox(slide, Inches(0.84), Inches(5.2), Inches(10.8), Inches(0.48), "Mood, focus, and review are not side features. They are the context that makes the AI useful.", size=18, color=WHITE, bold=True)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_kicker(slide, "AI FEATURES")
    add_title(slide, "Where AI is actually used in Daymark")
    add_card(slide, Inches(0.8), Inches(2.05), Inches(3.8), Inches(1.8), "1. Planning intelligence", "The planning engine ranks tasks, protects non-negotiables, proposes deferrals, and actively questions what should be deleted.", CYAN)
    add_card(slide, Inches(4.85), Inches(2.05), Inches(3.8), Inches(1.8), "2. Concierge reasoning", "The in-app concierge answers schedule questions, resolves follow-up references, and executes create / update / delete / defer / complete actions.", TEAL)
    add_card(slide, Inches(0.8), Inches(4.15), Inches(3.8), Inches(1.8), "3. Review summaries", "Daily, weekly, and monthly reviews are generated from actual product traces: task history, mood, focus, and future arrangements.", BLUE)
    add_card(slide, Inches(4.85), Inches(4.15), Inches(3.8), Inches(1.8), "4. Contextual generation", "The model also powers music recommendations and a daily fortune card using the user's current task load, mood, and context.", LIME)
    add_box(slide, Inches(9.05), Inches(2.2), Inches(3.15), Inches(3.65), SOFT_PANEL, alpha=0.08, line=RGBColor(80, 95, 118))
    add_textbox(slide, Inches(9.35), Inches(2.55), Inches(2.6), Inches(0.55), "Grounding principle", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.32), Inches(3.15), Inches(2.65), Inches(1.75), "Every AI feature is grounded in real app state: tasks, dates, mood, focus, review history, and profile data.", size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_background(slide, bright=True)
    add_kicker(slide, "MODEL + SYSTEM DESIGN")
    add_title(slide, "The core model strategy is LLM-first with reliable fallbacks.")
    add_card(slide, Inches(0.78), Inches(2.2), Inches(3.9), Inches(2.8), "Primary model", "DeepSeek Chat via an OpenAI-compatible chat completions API.\n\nRuntime configurable in Settings.\nBilingual prompts in English and Chinese.\nUsed for planning, concierge, reviews, songs, and fortune.", CYAN)
    add_textbox(slide, Inches(5.05), Inches(2.15), Inches(6.7), Inches(0.4), "Architecture flow", size=20, color=WHITE, bold=True)
    flow_x = Inches(5.05)
    flow = [
        ("App state", CYAN),
        ("Context builder", BLUE),
        ("Structured prompt", TEAL),
        ("DeepSeek", LIME),
        ("Validated output", ORANGE),
    ]
    for label, accent in flow:
        add_pipeline_step(slide, flow_x, Inches(2.75), Inches(1.18), Inches(0.58), label, accent)
        flow_x += Inches(1.26)
    add_textbox(slide, Inches(5.08), Inches(3.75), Inches(6.7), Inches(1.2), "If the model is unavailable, the system falls back to deterministic local behavior so the product remains usable and testable.", size=17, color=MUTED)
    add_textbox(slide, Inches(5.08), Inches(4.95), Inches(6.75), Inches(0.95), "Important product choice: model output is not blindly trusted. The backend validates structure and executes only allowed actions.", size=18, color=WHITE, bold=True)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_kicker(slide, "HOW I APPLY AI")
    add_title(slide, "I use AI as a reasoning layer, not just a chat layer.")
    add_card(slide, Inches(0.82), Inches(2.05), Inches(3.55), Inches(3.2), "Grounding layer", "Inputs include active tasks, due dates, recurring structure, mood logs, focus sessions, review history, birthday, and current date context.", CYAN)
    add_card(slide, Inches(4.6), Inches(2.05), Inches(3.55), Inches(3.2), "Reasoning layer", "The model parses intent, selects tasks, explains trade-offs, writes concise summaries, and generates context-aware content.", BLUE)
    add_card(slide, Inches(8.38), Inches(2.05), Inches(3.2), Inches(3.2), "Execution layer", "The product converts model output into task mutations, Today plans, review text, and recommendation payloads that the UI can safely render.", LIME)
    add_textbox(slide, Inches(0.86), Inches(5.7), Inches(10.8), Inches(0.55), "This is the central technical idea of the project: AI is most valuable when it is grounded, structured, and tied to real product actions.", size=19, color=WHITE, bold=True)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_background(slide, bright=True)
    add_kicker(slide, "ENGINEERING BUILD")
    add_title(slide, "What is inside the project")
    add_card(slide, Inches(0.82), Inches(2.15), Inches(3.3), Inches(2.9), "Frontend", "React 18\nReview calendar\nMood timeline\nFocus timer\nPrivate concierge\nBilingual UI", CYAN)
    add_card(slide, Inches(4.35), Inches(2.15), Inches(3.3), Inches(2.9), "Backend", "FastAPI + Pydantic\nTask / plan / review routers\nAssistant orchestration\nLLM runtime settings\nProtected demo seeding", TEAL)
    add_card(slide, Inches(7.88), Inches(2.15), Inches(3.35), Inches(2.9), "Data + reliability", "SQLAlchemy + SQLite\nTask history\nMood + focus traces\nReview snapshots\nLocal fallback paths\nRegression tests", LIME)
    add_textbox(slide, Inches(0.86), Inches(5.55), Inches(10.8), Inches(0.75), "The system is designed so that every important AI surface is backed by product state, persistence, and guardrails.", size=18, color=WHITE, bold=True)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_kicker(slide, "WHY THIS FITS ANOTE")
    add_title(slide, "This project follows the same human-centered AI philosophy.")
    add_box(slide, Inches(0.9), Inches(2.3), Inches(10.7), Inches(0.85), LIME, alpha=0.0)
    add_textbox(slide, Inches(1.15), Inches(2.52), Inches(10.1), Inches(0.45), "AI should help people make better decisions, not just produce more output.", size=26, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(1.05), Inches(3.6), Inches(10.2), [
        "It stays close to real human signals: time, energy, mood, focus, and trade-offs.",
        "It makes AI accessible in an everyday personal workflow, not only in expert or enterprise settings.",
        "It shows how model intelligence becomes more useful when embedded inside a product loop.",
    ], size=20, color=WHITE, bullet_color=CYAN, gap=0.72)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_background(slide, bright=True)
    add_kicker(slide, "NEXT STEPS")
    add_title(slide, "What I would improve next")
    add_card(slide, Inches(0.82), Inches(2.15), Inches(2.85), Inches(2.45), "Personalization", "Learn better daily capacity from real behavior rather than a fixed starting point.", CYAN)
    add_card(slide, Inches(3.95), Inches(2.15), Inches(2.85), Inches(2.45), "Evaluation", "Add stronger benchmarks for assistant reliability and planning quality.", BLUE)
    add_card(slide, Inches(7.08), Inches(2.15), Inches(2.85), Inches(2.45), "Longer horizon", "Expand from today-level reasoning into multi-day and weekly planning.", TEAL)
    add_card(slide, Inches(10.2), Inches(2.15), Inches(2.0), Inches(2.45), "Agents", "Evolve the concierge into a more robust task agent with measurable safety.", LIME)
    add_textbox(slide, Inches(0.86), Inches(5.45), Inches(10.5), Inches(0.85), "Daymark is my attempt to bring AI from \"interesting\" to \"useful\" at the scale of one real human day.", size=22, color=WHITE, bold=True)

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(str(OUTPUT_PATH))
    print(f"Saved presentation to: {OUTPUT_PATH}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
